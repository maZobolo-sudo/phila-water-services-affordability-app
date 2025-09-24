
import os
from datetime import datetime
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm
from pptx import Presentation
from pptx.util import Inches

def make_pdf_report(rows: list, path):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    c = canvas.Canvas(path, pagesize=A4); w,h=A4
    c.setFont("Helvetica-Bold",16); c.drawString(2*cm,h-2*cm,"WSA Affordability Summary")
    c.setFont("Helvetica",11); y=h-3*cm
    c.drawString(2*cm,y,f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"); y-=1*cm
    for r in rows[:25]:
        line = f"{r.get('wsa_name','')} | indigent {r.get('indigent_share',0):.0%} | FBW {r.get('fbw_kl',0)} kL | aff_typical_median {r.get('aff_typical_median',0):.3f}"
        c.drawString(2*cm,y,line); y-=0.7*cm
        if y<3*cm: c.showPage(); y=h-3*cm; c.setFont("Helvetica",11)
    c.showPage(); c.save(); return path

def make_pptx(rows: list, path):
    prs = Presentation(); s=prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text="WSA Affordability – Results"
    s.placeholders[1].text="Auto-generated summary"
    s2=prs.slides.add_slide(prs.slide_layouts[5]); tx=s2.shapes.add_textbox(Inches(1),Inches(1),Inches(8),Inches(5)).text_frame
    tx.text="Key Rows (first 15)"
    for r in rows[:15]:
        p = tx.add_paragraph()
        p.text = f"{r.get('wsa_name','')}: aff_typical_median={r.get('aff_typical_median',0):.3f} FBW={r.get('fbw_kl',0)}kL indigent={r.get('indigent_share',0):.0%}"
        p.level = 1
    prs.save(path); return path
